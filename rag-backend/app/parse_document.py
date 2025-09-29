# app/parse_document.py
# FastAPI Router: Parse Document (Sync + Async + Jobs + Env)
# Unterstützte Typen: PDF, DOCX, ODT*, XLSX, CSV, HTML, TXT, PPTX, EML
# - Upload (multipart) ODER file_url (async)
# - Einheitliche Ausgabe: { text, sections[], meta{} }
# - Schlanke, robuste Parser mit optionalen Abhängigkeiten
# - Keine LLMs; nur Extraktion/Normalisierung

import os
import re
import io
import uuid
import json
import time
import shutil
import string
import logging
import tempfile
from typing import Optional, List, Dict, Any, Tuple

from fastapi import APIRouter, UploadFile, File, Form, BackgroundTasks, HTTPException
from fastapi.responses import JSONResponse
from pydantic import BaseModel, AnyHttpUrl

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/parse", tags=["documents"])

# ------------- ENV / Defaults -------------
JOBS_DIR         = os.getenv("JOBS_DIR", "/data/jobs")
os.makedirs(JOBS_DIR, exist_ok=True)

# Optional: maximale Länge für text (0=unbegrenzt)
MAX_TEXT_CHARS   = int(os.getenv("PARSE_MAX_TEXT_CHARS", "0"))

# ------------- Utilities ------------------
def _job_path(jid: str) -> str:
    return os.path.join(JOBS_DIR, f"{jid}.parse.json")

def _job_save(j: Dict[str, Any]) -> None:
    with open(_job_path(j["id"]), "w", encoding="utf-8") as f:
        json.dump(j, f, ensure_ascii=False)

def _job_load(jid: str) -> Dict[str, Any]:
    p = _job_path(jid)
    if not os.path.exists(p):
        raise HTTPException(404, "job not found")
    with open(p, "r", encoding="utf-8") as f:
        return json.load(f)

def _download_to_tmp(url: str, prefix: str = "parse_") -> str:
    import requests
    tmpdir = tempfile.mkdtemp(prefix=prefix)
    dst = os.path.join(tmpdir, "input.bin")
    with requests.get(url, stream=True, timeout=60) as r:
        r.raise_for_status()
        with open(dst, "wb") as f:
            for chunk in r.iter_content(1024 * 1024):
                if chunk:
                    f.write(chunk)
    return dst

def _sniff_ext_mime(path: str, provided_mime: Optional[str] = None) -> Tuple[str, str]:
    mime = provided_mime or ""
    ext = os.path.splitext(path)[1].lower()
    if not mime:
        # best effort
        mapping = {
            ".pdf": "application/pdf",
            ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            ".odt": "application/vnd.oasis.opendocument.text",
            ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            ".csv": "text/csv",
            ".html": "text/html",
            ".htm": "text/html",
            ".txt": "text/plain",
            ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ".eml": "message/rfc822",
        }
        mime = mapping.get(ext, "application/octet-stream")
    return ext, mime

def _normalize_text(s: str) -> str:
    s = s.replace("\r\n", "\n").replace("\r", "\n")
    # collapse 3+ newlines to 2
    s = re.sub(r"\n{3,}", "\n\n", s)
    # trim trailing spaces per line
    s = "\n".join(line.rstrip() for line in s.split("\n"))
    return s.strip()

def _truncate_text(s: str) -> str:
    if MAX_TEXT_CHARS and len(s) > MAX_TEXT_CHARS:
        return s[:MAX_TEXT_CHARS]
    return s

# ------------- Parsers ---------------------
def _parse_pdf(path: str) -> Tuple[str, List[dict], Dict[str, Any]]:
    """PDF → Text per Seite + gesamter Text. Benötigt pdfminer.six."""
    try:
        from pdfminer.high_level import extract_pages
        from pdfminer.layout import LTTextContainer, LTTextBox, LTTextLine
    except Exception as e:
        raise HTTPException(503, detail=f"pdfminer.six not available: {e}")
    pages = []
    full = []
    page_no = 0
    for page_layout in extract_pages(path):
        page_no += 1
        buf = []
        for element in page_layout:
            if isinstance(element, (LTTextContainer, LTTextBox, LTTextLine)):
                txt = element.get_text() or ""
                if txt.strip():
                    buf.append(txt)
        page_text = _normalize_text("\n".join(buf))
        pages.append({"type": "page", "page": page_no, "text": page_text})
        full.append(page_text)
    text = _truncate_text("\n\n".join(full))
    meta = {"pages": page_no}
    return text, pages, meta

def _parse_docx(path: str) -> Tuple[str, List[dict], Dict[str, Any]]:
    try:
        import docx
    except Exception as e:
        raise HTTPException(503, detail=f"python-docx not available: {e}")
    doc = docx.Document(path)
    paras = []
    sections = []
    for i, p in enumerate(doc.paragraphs, start=1):
        t = (p.text or "").strip()
        if t:
            paras.append(t)
            sections.append({"type": "paragraph", "index": i, "text": t})
    text = _truncate_text(_normalize_text("\n\n".join(paras)))
    meta = {"paragraphs": len(paras)}
    return text, sections, meta

def _parse_odt(path: str) -> Tuple[str, List[dict], Dict[str, Any]]:
    try:
        from odf import text, teletype
        from odf.opendocument import load
    except Exception as e:
        raise HTTPException(503, detail=f"odfpy not available: {e}")
    doc = load(path)
    ps = doc.getElementsByType(text.P)
    paras = [teletype.extractText(p) for p in ps if teletype.extractText(p).strip()]
    sections = [{"type": "paragraph", "index": i+1, "text": p} for i, p in enumerate(paras)]
    text_out = _truncate_text(_normalize_text("\n\n".join(paras)))
    meta = {"paragraphs": len(paras)}
    return text_out, sections, meta

def _parse_xlsx(path: str) -> Tuple[str, List[dict], Dict[str, Any]]:
    try:
        import pandas as pd
    except Exception as e:
        raise HTTPException(503, detail=f"pandas not available: {e}")
    xls = pd.ExcelFile(path)
    texts = []
    sections: List[dict] = []
    for sheet in xls.sheet_names:
        df = xls.parse(sheet_name=sheet, dtype=str).fillna("")
        # in CSV-ähnlichen Text umwandeln
        csv_lines = [",".join(map(str, df.columns))]
        for _, row in df.iterrows():
            csv_lines.append(",".join(map(lambda x: str(x).replace("\n", " "), row.tolist())))
        t = _normalize_text("\n".join(csv_lines))
        sections.append({"type": "sheet", "sheet": sheet, "rows": int(df.shape[0]), "cols": int(df.shape[1]), "text": t})
        texts.append(f"# {sheet}\n{t}")
    text = _truncate_text("\n\n".join(texts))
    meta = {"sheets": len(xls.sheet_names), "sheet_names": xls.sheet_names}
    return text, sections, meta

def _parse_csv(path: str) -> Tuple[str, List[dict], Dict[str, Any]]:
    try:
        import pandas as pd
    except Exception as e:
        raise HTTPException(503, detail=f"pandas not available: {e}")
    df = pd.read_csv(path, dtype=str).fillna("")
    csv_lines = [",".join(map(str, df.columns))]
    for _, row in df.iterrows():
        csv_lines.append(",".join(map(lambda x: str(x).replace("\n", " "), row.tolist())))
    t = _normalize_text("\n".join(csv_lines))
    sections = [{"type": "csv", "rows": int(df.shape[0]), "cols": int(df.shape[1]), "text": t}]
    text = _truncate_text(t)
    meta = {"rows": int(df.shape[0]), "cols": int(df.shape[1])}
    return text, sections, meta

def _parse_html(path: str) -> Tuple[str, List[dict], Dict[str, Any]]:
    try:
        from bs4 import BeautifulSoup
    except Exception as e:
        raise HTTPException(503, detail=f"beautifulsoup4 not available: {e}")
    with open(path, "rb") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    # Headlines + Paragraphs als Abschnitte
    sections = []
    texts = []
    for h in soup.find_all(re.compile("^h[1-6]$")):
        title = h.get_text(" ", strip=True)
        block = [title]
        for sib in h.find_all_next():
            if sib.name and re.match("^h[1-6]$", sib.name, re.I):
                break
            if sib.name == "p":
                txt = sib.get_text(" ", strip=True)
                if txt:
                    block.append(txt)
        sec_text = _normalize_text("\n".join(block))
        if sec_text:
            sections.append({"type": "html_section", "heading": title, "text": sec_text})
            texts.append(sec_text)
    if not texts:
        # Fallback: gesamter sichtbarer Text
        txt = soup.get_text(" ", strip=True)
        texts = [_normalize_text(txt)]
        sections = [{"type": "html", "text": texts[0]}]
    full = _truncate_text("\n\n".join(texts))
    meta = {"sections": len(sections)}
    return full, sections, meta

def _parse_txt(path: str) -> Tuple[str, List[dict], Dict[str, Any]]:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        raw = f.read()
    text = _truncate_text(_normalize_text(raw))
    sections = [{"type": "text", "text": text}]
    meta = {"chars": len(text)}
    return text, sections, meta

def _parse_pptx(path: str) -> Tuple[str, List[dict], Dict[str, Any]]:
    try:
        from pptx import Presentation
    except Exception as e:
        raise HTTPException(503, detail=f"python-pptx not available: {e}")
    prs = Presentation(path)
    sections = []
    texts = []
    for idx, slide in enumerate(prs.slides, start=1):
        buf = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    buf.append(t)
        st = _normalize_text("\n".join(buf))
        sections.append({"type": "slide", "slide": idx, "text": st})
        if st:
            texts.append(f"# Folie {idx}\n{st}")
    text = _truncate_text("\n\n".join(texts))
    meta = {"slides": len(prs.slides)}
    return text, sections, meta

def _parse_eml(path: str) -> Tuple[str, List[dict], Dict[str, Any]]:
    import email
    from email import policy
    with open(path, "rb") as f:
        msg = email.message_from_binary_file(f, policy=policy.default)
    subject = msg.get("subject", "")
    parts = []
    for part in msg.walk():
        if part.get_content_type() in ("text/plain", "text/html"):
            try:
                payload = part.get_payload(decode=True) or b""
                if not payload:
                    continue
                if part.get_content_type() == "text/html":
                    try:
                        from bs4 import BeautifulSoup
                        txt = BeautifulSoup(payload, "html.parser").get_text(" ", strip=True)
                    except Exception:
                        txt = payload.decode(errors="ignore")
                else:
                    txt = payload.decode(errors="ignore")
                if txt:
                    parts.append(txt)
            except Exception:
                continue
    body = _normalize_text("\n\n".join(parts))
    text = _truncate_text((subject + "\n\n" + body).strip())
    sections = [{"type": "email", "subject": subject, "text": text}]
    meta = {"subject": subject}
    return text, sections, meta

# Dispatcher
def _parse_by_ext(path: str, ext: str) -> Tuple[str, List[dict], Dict[str, Any]]:
    if ext == ".pdf":
        return _parse_pdf(path)
    if ext == ".docx":
        return _parse_docx(path)
    if ext == ".odt":
        return _parse_odt(path)
    if ext == ".xlsx":
        return _parse_xlsx(path)
    if ext == ".csv":
        return _parse_csv(path)
    if ext in (".html", ".htm"):
        return _parse_html(path)
    if ext == ".txt":
        return _parse_txt(path)
    if ext == ".pptx":
        return _parse_pptx(path)
    if ext == ".eml":
        return _parse_eml(path)
    # Fallback: als Text lesen
    return _parse_txt(path)

# ------------- Models ----------------------
class ParseOut(BaseModel):
    text: str
    sections: List[dict]
    meta: dict
    debug: Optional[dict] = None

class ParseAsyncIn(BaseModel):
    file_url: AnyHttpUrl
    mime: Optional[str] = None
    callback_url: Optional[AnyHttpUrl] = None
    meta: Optional[dict] = None

# ------------- SYNC ------------------------
@router.post("/document", response_model=ParseOut)
async def parse_document(
    file: UploadFile = File(...),
    mime: Optional[str] = Form(default=None)
):
    """
    Dokument-Parsing (synchron) per Upload.
    Gibt text + sections + meta zurück.
    """
    tmpdir = tempfile.mkdtemp(prefix="parse_")
    src_path = os.path.join(tmpdir, f"src_{uuid.uuid4().hex}")
    try:
        # persist upload
        with open(src_path, "wb") as f:
            while True:
                chunk = await file.read(1024 * 1024)
                if not chunk:
                    break
                f.write(chunk)

        ext, detected_mime = _sniff_ext_mime(src_path if file.filename is None else file.filename, mime or file.content_type)
        text, sections, meta = _parse_by_ext(src_path, ext)
        out = {
            "text": text,
            "sections": sections,
            "meta": {
                **meta,
                "filename": file.filename,
                "mime": detected_mime,
                "ext": ext,
                "chars": len(text)
            },
            "debug": { "workdir": os.path.basename(tmpdir) }
        }
        return JSONResponse(out)
    except HTTPException:
        raise
    except Exception as e:
        logger.exception("parse_document failed")
        return JSONResponse({"error": str(e)}, status_code=500)
    finally:
        try:
            shutil.rmtree(tmpdir, ignore_errors=True)
        except Exception:
            pass

# ------------- ASYNC + JOBS ----------------
@router.post("/document/async")
def parse_document_async(body: ParseAsyncIn, bg: BackgroundTasks):
    jid = uuid.uuid4().hex
    job = {
        "id": jid,
        "status": "queued",
        "created_at": int(time.time()),
        "updated_at": int(time.time()),
        "request": body.dict(),
        "result": None,
        "error": None,
    }
    _job_save(job)
    bg.add_task(_do_parse_job, jid)
    return {"job_id": jid, "status": "queued"}

@router.get("/jobs/{job_id}")
def parse_job_status(job_id: str):
    return _job_load(job_id)

def _post_callback(url: str, payload: dict) -> None:
    try:
        import requests
        requests.post(str(url), json=payload, timeout=15)
    except Exception:
        logger.warning("callback post failed", exc_info=True)

def _do_parse_job(job_id: str) -> None:
    try:
        j = _job_load(job_id)
        req = j.get("request") or {}
        src = _download_to_tmp(req["file_url"])
        filename = req.get("file_url")
        ext, detected_mime = _sniff_ext_mime(filename or src, req.get("mime"))

        text, sections, meta = _parse_by_ext(src, ext)
        result = {
            "text": text,
            "sections": sections,
            "meta": { **meta, "filename": filename, "mime": detected_mime, "ext": ext, "chars": len(text) },
            "debug": { "source_url": filename }
        }

        j.update({"status": "done", "result": result, "updated_at": int(time.time())})
        _job_save(j)

        if req.get("callback_url"):
            _post_callback(req["callback_url"], {"job_id": j["id"], "status": "done", "result": result})
        try:
            shutil.rmtree(os.path.dirname(src), ignore_errors=True)
        except Exception:
            pass

    except Exception as e:
        logger.exception("parse job failed")
        try:
            j = _job_load(job_id)
            j.update({"status": "error", "error": str(e), "updated_at": int(time.time())})
            _job_save(j)
            req = j.get("request") or {}
            if req.get("callback_url"):
                _post_callback(req["callback_url"], {"job_id": j["id"], "status": "error", "error": str(e)})
        except Exception:
            pass

# ------------- ENV Debug -------------------
@router.get("/env")
def show_env():
    libs = {}
    for name in ("pdfminer", "docx", "odf", "pandas", "bs4", "pptx"):
        try:
            __import__(name)
            libs[name] = True
        except Exception:
            libs[name] = False
    return {
        "JOBS_DIR": JOBS_DIR,
        "PARSE_MAX_TEXT_CHARS": MAX_TEXT_CHARS,
        "available_libs": libs
    }
